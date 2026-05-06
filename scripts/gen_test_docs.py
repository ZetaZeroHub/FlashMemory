#!/usr/bin/env python3
"""
Generate test fixtures for DocParser unit tests.

Dependencies:
    pip install python-docx python-pptx Pillow reportlab

Optional:
    brew install pandoc       # used by some tests; not required for fixture generation

Output: internal/parser/docs/testdata/{md,pdf,docx,pptx,images}/
"""

import os
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
TESTDATA = ROOT / "internal" / "parser" / "docs" / "testdata"


# ---------- Markdown ----------

def gen_md_fixtures():
    out = TESTDATA / "md"
    out.mkdir(parents=True, exist_ok=True)

    (out / "kb_sample.md").write_text(
        "# 测试知识库\n\n"
        "## 信用卡业务\n\n"
        "### 1. 我即将出国旅行,能否帮我开通信用卡境外使用功能?\n\n"
        "您好,可以为您开通信用卡境外交易功能。\n\n"
        "*[条目编号:`kb_qa_98a5cb04ff`]*\n\n"
        "---\n\n"
        "### 2. 我想查看信用卡的年费,该怎么操作?\n\n"
        "登录手机银行APP查看。\n\n"
        "*[条目编号:`kb_qa_c03f5681a2`]*\n\n"
        "---\n\n"
        "### 3. 我想注销信用卡。\n\n"
        "您可在线申请注销。\n\n"
        "*[条目编号:`kb_qa_3a1d5f2d55`]*\n",
        encoding="utf-8",
    )

    (out / "nested_headings.md").write_text(
        "# 顶级\n\n顶级正文\n\n## 二级 A\n\n二级 A 正文\n\n### 三级 A.1\n\n"
        "三级 A.1 正文\n\n### 三级 A.2\n\n三级 A.2 正文\n\n## 二级 B\n\n二级 B 正文\n",
        encoding="utf-8",
    )

    (out / "no_headings.md").write_text(
        "\n".join([f"line {i} 内容" for i in range(1, 200)]),
        encoding="utf-8",
    )

    long_body = "段落 " + ("中文测试 " * 1500)  # > 5000 chars
    (out / "long_section.md").write_text(
        f"# 长文档\n\n## 唯一章节\n\n{long_body}\n\n*[条目编号:`kb_qa_aaaa1111ff`]*\n",
        encoding="utf-8",
    )

    (out / "intro_before_heading.md").write_text(
        "这是引言部分,在第一个标题之前。\n\n# 第一个标题\n\n标题正文\n",
        encoding="utf-8",
    )

    (out / "emoji_titles.md").write_text(
        "# 🚀 Rocket标题\n\n正文\n\n## 标题 with @special! chars#\n\n正文\n",
        encoding="utf-8",
    )

    print(f"  [md] 6 fixtures written to {out}")


# ---------- DOCX ----------

def gen_docx_fixtures():
    try:
        from docx import Document
        from docx.shared import Pt
    except ImportError:
        print("  [docx] SKIP (python-docx not installed)")
        return

    out = TESTDATA / "docx"
    out.mkdir(parents=True, exist_ok=True)

    # 英文 Heading 样式
    doc = Document()
    doc.add_heading("Chapter One", level=1)
    doc.add_paragraph("Chapter one body text.")
    doc.add_heading("Section 1.1", level=2)
    doc.add_paragraph("Section 1.1 body.")
    doc.add_heading("Subsection 1.1.1", level=3)
    doc.add_paragraph("Subsection content.")
    doc.add_heading("Chapter Two", level=1)
    doc.add_paragraph("Chapter two body. *[条目编号:`kb_qa_bbbb2222ff`]*")
    doc.save(out / "headings_en.docx")

    # 中文版样式（python-docx 默认 Heading 1，但我们手动设置 style 名）
    doc = Document()
    p1 = doc.add_paragraph("第一章")
    p1.style = doc.styles["Heading 1"]
    doc.add_paragraph("第一章正文。 *[条目编号:`kb_qa_cccc3333ff`]*")
    p2 = doc.add_paragraph("第一节")
    p2.style = doc.styles["Heading 2"]
    doc.add_paragraph("第一节正文。")
    doc.save(out / "headings_cn.docx")

    # 含表格
    doc = Document()
    doc.add_heading("含表格的文档", level=1)
    doc.add_paragraph("表格前的段落。")
    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "列A"
    table.rows[0].cells[1].text = "列B"
    table.rows[0].cells[2].text = "列C"
    table.rows[1].cells[0].text = "数据1"
    table.rows[1].cells[1].text = "数据2"
    table.rows[1].cells[2].text = "数据3"
    doc.add_paragraph("表格后的段落。 *[条目编号:`kb_qa_dddd4444ff`]*")
    doc.save(out / "with_table.docx")

    # 无样式
    doc = Document()
    for i in range(40):
        doc.add_paragraph(f"段落 {i+1} 的文本内容。")
    doc.save(out / "no_styles.docx")

    print(f"  [docx] 4 fixtures written to {out}")


# ---------- PPTX ----------

def gen_pptx_fixtures():
    try:
        from pptx import Presentation
    except ImportError:
        print("  [pptx] SKIP (python-pptx not installed)")
        return

    out = TESTDATA / "pptx"
    out.mkdir(parents=True, exist_ok=True)

    # ≥10 张 slide 测排序
    prs = Presentation()
    blank_layout = prs.slide_layouts[5]
    for i in range(1, 12):
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = f"Slide {i}"
        body = slide.shapes.add_textbox(left=914400, top=2057400, width=7315200, height=4572000)
        body.text_frame.text = f"This is slide number {i} content."
    prs.save(out / "multi_slides.pptx")

    # 含备注
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide With Notes"
    notes = slide.notes_slide.notes_text_frame
    notes.text = "This is the speaker note. Key talking point."
    prs.save(out / "with_notes.pptx")

    print(f"  [pptx] 2 fixtures written to {out}")


# ---------- PDF ----------

def gen_pdf_fixtures():
    try:
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    except ImportError:
        print("  [pdf] SKIP (reportlab not installed)")
        return

    out = TESTDATA / "pdf"
    out.mkdir(parents=True, exist_ok=True)

    try:
        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        cn_font = "STSong-Light"
    except Exception:
        cn_font = None

    # 英文数字标题
    c = canvas.Canvas(str(out / "english_numeric.pdf"), pagesize=A4)
    c.setFont("Helvetica", 12)
    y = 800
    for line in [
        "1. INTRODUCTION",
        "This is the intro.",
        "1.1 Background",
        "Background details.",
        "1.2 Scope",
        "Scope details.",
        "2. METHODOLOGY",
        "Method details. [条目编号:kb_qa_eeee5555ff]",
    ]:
        c.drawString(60, y, line)
        y -= 20
    c.save()

    # 中文标题
    if cn_font:
        c = canvas.Canvas(str(out / "chinese_headings.pdf"), pagesize=A4)
        c.setFont(cn_font, 12)
        y = 800
        for line in [
            "第一章 总则",
            "总则正文。",
            "1. 我想申诉话费",
            "请按以下步骤操作。 [条目编号:kb_qa_ffff6666ff]",
            "2. 我想查询余额",
            "登录手机银行查询。 [条目编号:kb_qa_gggg7777ff]",
            "第二章 业务办理",
            "业务办理内容。",
        ]:
            c.drawString(60, y, line)
            y -= 20
        c.save()
    else:
        print("  [pdf] WARN: Chinese font unavailable, skipping chinese_headings.pdf")

    # 跨页 QA
    c = canvas.Canvas(str(out / "cross_page_qa.pdf"), pagesize=A4)
    c.setFont("Helvetica", 12)
    # 第一页：QA 开头
    c.drawString(60, 800, "Question: long question that spans pages")
    for i in range(35):
        c.drawString(60, 780 - i * 20, f"Line {i+1} of answer text.")
    c.showPage()
    # 第二页：QA 结尾
    c.setFont("Helvetica", 12)
    c.drawString(60, 800, "Continued answer text on page 2.")
    c.drawString(60, 780, "Final line. [条目编号:kb_qa_hhhh8888ff]")
    c.save()

    # 单页平铺无标题
    c = canvas.Canvas(str(out / "flat_page.pdf"), pagesize=A4)
    c.setFont("Helvetica", 10)
    for i in range(80):
        c.drawString(60, 800 - i * 9, f"flat line content number {i+1}")
    c.save()

    print(f"  [pdf] fixtures written to {out}")


# ---------- Image ----------

def gen_image_fixtures():
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("  [images] SKIP (Pillow not installed)")
        return

    out = TESTDATA / "images"
    out.mkdir(parents=True, exist_ok=True)

    def make_image(name, text, size=(400, 100)):
        img = Image.new("RGB", size, "white")
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("/System/Library/Fonts/PingFang.ttc", 24)
        except Exception:
            font = ImageFont.load_default()
        draw.text((10, 30), text, fill="black", font=font)
        img.save(out / name)

    make_image("chinese_text.png", "信用卡境外开通")
    make_image("english_text.png", "Hello World OCR")
    make_image("low_quality.png", "blur test", size=(120, 40))

    print(f"  [images] 3 fixtures written to {out}")


def main():
    print(f"Generating test fixtures into: {TESTDATA}")
    TESTDATA.mkdir(parents=True, exist_ok=True)
    gen_md_fixtures()
    gen_docx_fixtures()
    gen_pptx_fixtures()
    gen_pdf_fixtures()
    gen_image_fixtures()
    print("Done.")


if __name__ == "__main__":
    main()
